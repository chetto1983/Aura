"""
Unit tests for AuraPptxConverter.

Run with:  python -m pytest docker/markitdown/aura-plugins/tests/test_pptx_converter.py -q
"""

import io
import os
import sys

# Allow running from repo root without installing the package
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from markitdown import StreamInfo
from pptx import Presentation
from pptx.util import Inches

from aura.pptx_converter import AuraPptxConverter


def _stream_info(filename: str = "deck.pptx") -> StreamInfo:
    return StreamInfo(
        extension=".pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
    )


def _save(prs: Presentation) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx() -> bytes:
    prs = Presentation()
    prs.core_properties.title = "Quarterly Review"
    prs.core_properties.author = "Aura Test"

    title_slide = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide)
    slide1.shapes.title.text = "Overview"
    slide1.placeholders[1].text = "Revenue grew in Q1."
    slide1.notes_slide.notes_text_frame.text = "Mention regional variance."

    title_only = prs.slide_layouts[5]
    slide2 = prs.slides.add_slide(title_only)
    slide2.shapes.title.text = "Risks"
    box = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    box.text_frame.text = "Supply chain delays need monitoring."

    return _save(prs)


def _make_blank_pptx() -> bytes:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    return _save(prs)


def test_accepts_pptx_extension():
    c = AuraPptxConverter()
    assert c.accepts(io.BytesIO(b""), _stream_info())


def test_accepts_pptx_mime():
    c = AuraPptxConverter()
    si = StreamInfo(
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert c.accepts(io.BytesIO(b""), si)


def test_rejects_docx():
    c = AuraPptxConverter()
    si = StreamInfo(
        extension=".docx",
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    assert not c.accepts(io.BytesIO(b""), si)


def test_convert_basic_structure_and_frontmatter():
    c = AuraPptxConverter()
    result = c.convert(io.BytesIO(_make_pptx()), _stream_info("review.pptx"))
    md = result.markdown

    assert md.startswith("---")
    assert 'presentation_title: "Quarterly Review"' in md
    assert 'author: "Aura Test"' in md
    assert "slide_count: 2" in md
    assert "slide_titles:" in md
    assert '  - "Overview"' in md
    assert '  - "Risks"' in md
    assert "has_notes: true" in md

    assert "### Slide 1: Overview" in md
    assert "### Slide 2: Risks" in md
    assert "Revenue grew in Q1." in md
    assert "Supply chain delays need monitoring." in md
    assert "**Speaker notes:** Mention regional variance." in md


def test_convert_untitled_slide_gets_placeholder():
    prs = Presentation()
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    box = blank.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.text = "Body without a title."

    c = AuraPptxConverter()
    result = c.convert(io.BytesIO(_save(prs)), _stream_info())

    assert '  - "(untitled)"' in result.markdown
    assert "### Slide 1: (untitled)" in result.markdown
    assert "Body without a title." in result.markdown


def test_zero_title_and_zero_text_deck_uses_builtin_fallback():
    c = AuraPptxConverter()
    result = c.convert(io.BytesIO(_make_blank_pptx()), _stream_info("blank.pptx"))

    assert not result.markdown.startswith("---")
    assert "<!-- Slide number: 1 -->" in result.markdown
    assert "### Slide 1:" not in result.markdown
