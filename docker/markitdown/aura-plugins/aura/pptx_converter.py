"""
Per-slide PPTX decomposition for Aura's wiki heading-subnode indexer.

Each slide becomes a ### Slide N heading so slide-level facts can be indexed
as separate retrieval nodes. Text-bearing shapes are emitted in package XML
order, which is usually visual order but is not guaranteed by python-pptx.
Speaker notes are appended when present.
"""

import io
import os
import warnings
from typing import Any, BinaryIO

from markitdown import DocumentConverter, DocumentConverterResult, StreamInfo
from markitdown.converters import PptxConverter
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

_PPTX_EXTENSIONS = {".pptx"}
_PPTX_MIME_PREFIXES = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
)

_TITLE_PLACEHOLDERS = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.VERTICAL_TITLE,
}


def _yaml_scalar(s: str) -> str:
    """Always-quoted YAML scalar; avoids parsing ambiguity for arbitrary strings."""
    safe = str(s).replace("\\", "\\\\").replace('"', '\\"').replace("\n", "\\n")
    return f'"{safe}"'


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return ""
    return (text_frame.text or "").strip()


def _is_title_shape(shape) -> bool:
    try:
        return shape.placeholder_format.type in _TITLE_PLACEHOLDERS
    except ValueError:
        return False


def _slide_title(slide) -> tuple[str, bool]:
    for shape in slide.shapes:
        if _is_title_shape(shape):
            text = _shape_text(shape)
            if text:
                return text, True

    return "(untitled)", False


def _notes_text(slide) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
    if notes_frame is None:
        return ""
    return (notes_frame.text or "").strip()


class AuraPptxConverter(DocumentConverter):
    """
    Per-slide PPTX converter. Output shape:
      ---
      presentation_title: "..."
      author: "..."
      slide_count: 2
      slide_titles:
        - "Intro"
        - "(untitled)"
      has_notes: true
      ---

      ### Slide 1: Intro
      Slide body text.

      **Speaker notes:** Presenter note.
    """

    def __init__(self) -> None:
        super().__init__()
        self._fallback = PptxConverter()

    def accepts(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,
    ) -> bool:
        ext = (getattr(stream_info, "extension", "") or "").lower()
        mime = (getattr(stream_info, "mimetype", "") or "").lower()
        return ext in _PPTX_EXTENSIONS or any(
            mime.startswith(p) for p in _PPTX_MIME_PREFIXES
        )

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,
    ) -> DocumentConverterResult:
        data = file_stream.read()

        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            prs = Presentation(io.BytesIO(data))

        slides = []
        title_shape_count = 0
        text_shape_count = 0
        has_notes = False

        for slide in prs.slides:
            title, title_from_shape = _slide_title(slide)
            if title_from_shape:
                title_shape_count += 1

            body_texts = []
            for shape in slide.shapes:
                text = _shape_text(shape)
                if not text:
                    continue
                text_shape_count += 1
                if title_from_shape and _is_title_shape(shape) and text == title:
                    continue
                body_texts.append(text)

            notes = _notes_text(slide)
            if notes:
                has_notes = True

            slides.append(
                {
                    "title": title,
                    "body_texts": body_texts,
                    "notes": notes,
                }
            )

        # If the deck has no extractable title or text, preserve markitdown's
        # built-in behavior rather than replacing it with empty slide headings.
        if title_shape_count == 0 and text_shape_count == 0:
            return self._fallback.convert(io.BytesIO(data), stream_info, **kwargs)

        props = prs.core_properties
        filename = (getattr(stream_info, "filename", None) or "").strip()
        fallback_title = os.path.splitext(os.path.basename(filename))[0] or "presentation"
        presentation_title = (props.title or "").strip() or fallback_title
        author = (props.author or "").strip() or "unknown"

        lines = ["---"]
        lines.append(f"presentation_title: {_yaml_scalar(presentation_title)}")
        lines.append(f"author: {_yaml_scalar(author)}")
        lines.append(f"slide_count: {len(slides)}")
        lines.append("slide_titles:")
        for slide in slides:
            lines.append(f"  - {_yaml_scalar(slide['title'])}")
        lines.append(f"has_notes: {str(has_notes).lower()}")
        lines.append("---")
        lines.append("")

        for idx, slide in enumerate(slides, start=1):
            lines.append(f"### Slide {idx}: {slide['title']}")
            lines.append("")
            for text in slide["body_texts"]:
                lines.append(text)
                lines.append("")
            if slide["notes"]:
                lines.append(f"**Speaker notes:** {slide['notes']}")
                lines.append("")

        return DocumentConverterResult(markdown="\n".join(lines).rstrip())
