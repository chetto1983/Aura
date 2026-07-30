# Aura document sidecar. It keeps /convert compatible with the original
# markitdown wrapper and adds /extract for structured document ingestion.
import base64
import os
import re
import shutil
import tempfile
from typing import Any

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from markitdown import MarkItDown

app = FastAPI(title="aura-markitdown", version="2")

# OCR consolidation: markitdown shares Aura's GLM-OCR engine (the aura-ocr-vl sidecar,
# an OpenAI-compatible vision endpoint) — the SAME engine the chat/vision path uses. When
# AURA_OCR_BASE_URL is set, standalone images and scanned-PDF pages are OCR'd through it
# and markitdown's own converters get the llm_client (so embedded images convert too).
# Fail-soft: no endpoint / unreachable -> text-only extraction exactly as before.
_ocr_base = os.getenv("AURA_OCR_BASE_URL", "").strip()
_ocr_model = os.getenv("AURA_OCR_MODEL", "glm-ocr").strip() or "glm-ocr"
_ocr_pdf_max_pages = int(os.getenv("AURA_OCR_PDF_MAX_PAGES", "40") or "40")
_ocr_timeout = float(os.getenv("AURA_OCR_TIMEOUT_SEC", "120") or "120")
_OCR_PROMPT = (
    "Extract ALL text from this image verbatim as clean Markdown. Preserve headings, "
    "tables, lists, and reading order. Output only the extracted content, no commentary."
)


def _build_ocr_client():
    if not _ocr_base:
        return None
    try:
        from openai import OpenAI

        return OpenAI(
            base_url=_ocr_base,
            api_key=os.getenv("AURA_OCR_API_KEY", "sk-noauth") or "sk-noauth",
            timeout=_ocr_timeout,
        )
    except Exception:
        return None


_ocr = _build_ocr_client()
if _ocr is not None:
    try:
        _md = MarkItDown(enable_plugins=True, llm_client=_ocr, llm_model=_ocr_model)
    except Exception:
        _md = MarkItDown()
else:
    _md = MarkItDown()

_image_mime = {
    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp",
    ".tif": "image/tiff", ".tiff": "image/tiff",
}
_ws_re = re.compile(r"\s+")
_atx_heading_re = re.compile(r"^(#{1,6})\s+(.*)$")
_max_chunk_chars = 12000


@app.get("/health")
def health() -> dict:
    return {"ok": True}


@app.post("/convert")
async def convert(file: UploadFile = File(...)) -> dict:
    path = _save_upload(file)
    try:
        result = _md.convert(path)
        return {"markdown": result.text_content}
    except Exception:
        raise HTTPException(status_code=422, detail="conversion failed")
    finally:
        _unlink(path)


@app.post("/extract")
async def extract(file: UploadFile = File(...), mime_type: str = Form("")) -> dict:
    path = _save_upload(file)
    suffix = os.path.splitext(file.filename or "")[1].lower()
    effective_mime = mime_type or file.content_type or ""
    try:
        if suffix == ".pdf" or effective_mime == "application/pdf":
            payload = _extract_pdf(path)
        elif suffix in (".xlsx", ".xlsm") or "spreadsheet" in effective_mime:
            payload = _extract_xlsx(path)
        elif suffix == ".docx" or effective_mime.endswith("wordprocessingml.document"):
            payload = _extract_docx(path)
        elif suffix == ".pptx" or effective_mime.endswith("presentationml.presentation"):
            payload = _extract_pptx(path)
        elif suffix in (".html", ".htm") or effective_mime in ("text/html", "application/xhtml+xml"):
            payload = _extract_html(path)
        elif suffix == ".csv" or effective_mime == "text/csv":
            payload = _extract_csv(path)
        elif suffix in _image_mime or effective_mime.startswith("image/"):
            payload = _extract_image(path, effective_mime)
        else:
            payload = _extract_markdown(path)
        payload["mime_type"] = effective_mime or payload.get("mime_type", "")
        return payload
    except HTTPException:
        raise
    except Exception:
        raise HTTPException(status_code=422, detail="extraction failed")
    finally:
        _unlink(path)


def _save_upload(file: UploadFile) -> str:
    suffix = os.path.splitext(file.filename or "")[1] or ".bin"
    fd, path = tempfile.mkstemp(suffix=suffix)
    try:
        with os.fdopen(fd, "wb") as fh:
            file.file.seek(0)
            shutil.copyfileobj(file.file, fh)
        return path
    except Exception:
        _unlink(path)
        raise


def _unlink(path: str) -> None:
    try:
        os.unlink(path)
    except FileNotFoundError:
        pass


def _normalize(text: str) -> str:
    return _ws_re.sub(" ", text or "").strip()


def _chunk_text(kind: str, text: str, locator: dict[str, Any], heading_path: list[str] | None = None) -> list[dict]:
    clean = _normalize(text)
    if not clean:
        return []
    heading_path = heading_path or []
    chunks = []
    while len(clean) > _max_chunk_chars:
        split_at = clean.rfind(" ", 0, _max_chunk_chars)
        if split_at < _max_chunk_chars // 2:
            split_at = _max_chunk_chars
        chunks.append(
            {
                "kind": kind,
                "text": clean[:split_at].strip(),
                "locator": locator,
                "heading_path": heading_path,
            }
        )
        clean = clean[split_at:].strip()
    if clean:
        chunks.append(
            {
                "kind": kind,
                "text": clean,
                "locator": locator,
                "heading_path": heading_path,
            }
        )
    return chunks


def _payload(title: str, mime_type: str, chunks: list[dict], pages: int = 0, sheets: int = 0, sections: int = 0) -> dict:
    return {
        "title": title,
        "mime_type": mime_type,
        "stats": {
            "pages": pages,
            "sheets": sheets,
            "sections": sections,
            "chunks": len(chunks),
            "characters": sum(len(chunk.get("text", "")) for chunk in chunks),
        },
        "chunks": chunks,
    }


def _ocr_image_bytes(data: bytes, mime: str) -> str:
    """OCR raw image bytes via the shared GLM-OCR engine. Returns '' on any failure (fail-soft)."""
    if _ocr is None or not data:
        return ""
    b64 = base64.b64encode(data).decode("ascii")
    try:
        resp = _ocr.chat.completions.create(
            model=_ocr_model,
            temperature=0,
            messages=[{"role": "user", "content": [
                {"type": "text", "text": _OCR_PROMPT},
                {"type": "image_url", "image_url": {"url": f"data:{mime or 'image/png'};base64,{b64}"}},
            ]}],
        )
        return (resp.choices[0].message.content or "").strip()
    except Exception:
        return ""


def _extract_image(path: str, mime: str) -> dict:
    suffix = os.path.splitext(path)[1].lower()
    eff_mime = mime or _image_mime.get(suffix, "image/png")
    with open(path, "rb") as fh:
        data = fh.read()
    if _ocr is None:
        raise HTTPException(status_code=503, detail="image OCR unavailable (no OCR engine configured)")
    text = _ocr_image_bytes(data, eff_mime)
    chunks = _chunk_text("image", text, {})
    return _payload(os.path.basename(path), eff_mime, chunks)


def _extract_pdf(path: str) -> dict:
    try:
        import fitz
    except Exception:
        raise HTTPException(status_code=503, detail="pdf extractor unavailable")

    doc = fitz.open(path)
    try:
        title = (doc.metadata or {}).get("title") or os.path.basename(path)
        chunks: list[dict] = []
        for index, page in enumerate(doc, start=1):
            text = page.get_text("text") or ""
            # Scanned/image-only page: render and OCR via the shared GLM-OCR engine (bounded).
            if not _normalize(text) and _ocr is not None and index <= _ocr_pdf_max_pages:
                try:
                    pix = page.get_pixmap(dpi=200)
                    text = _ocr_image_bytes(pix.tobytes("png"), "image/png")
                except Exception:
                    text = ""
            chunks.extend(_chunk_text("page", text, {"page": index}))
        return _payload(title, "application/pdf", chunks, pages=doc.page_count)
    finally:
        doc.close()


def _extract_xlsx(path: str) -> dict:
    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
    except Exception:
        raise HTTPException(status_code=503, detail="xlsx extractor unavailable")

    wb = load_workbook(path, read_only=True, data_only=False)
    try:
        chunks: list[dict] = []
        for ws in wb.worksheets:
            lines: list[str] = []
            row_start = 0
            row_end = 0
            for row_index, row in enumerate(ws.iter_rows(), start=1):
                cells = []
                for column_index, cell in enumerate(row, start=1):
                    value = cell.value
                    if value is None:
                        continue
                    text = str(value).strip()
                    if not text:
                        continue
                    cells.append(f"{get_column_letter(column_index)}{row_index}: {text}")
                if not cells:
                    continue
                if row_start == 0:
                    row_start = row_index
                row_end = row_index
                lines.append(f"row {row_index}: " + " | ".join(cells))
                joined = "\n".join(lines)
                if len(joined) >= _max_chunk_chars or len(lines) >= 50:
                    chunks.extend(
                        _chunk_text(
                            "rows",
                            joined,
                            {"sheet": ws.title, "row_start": row_start, "row_end": row_end},
                            [ws.title],
                        )
                    )
                    lines = []
                    row_start = 0
                    row_end = 0
            if lines:
                chunks.extend(
                    _chunk_text(
                        "rows",
                        "\n".join(lines),
                        {"sheet": ws.title, "row_start": row_start, "row_end": row_end},
                        [ws.title],
                    )
                )
        return _payload(os.path.basename(path), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", chunks, sheets=len(wb.worksheets))
    finally:
        wb.close()


def _extract_docx(path: str) -> dict:
    try:
        from docx import Document
    except Exception:
        raise HTTPException(status_code=503, detail="docx extractor unavailable")

    doc = Document(path)
    chunks: list[dict] = []
    heading_path: list[str] = []
    buffer: list[str] = []
    start_para = 0
    sections = 0

    def flush(end_para: int) -> None:
        nonlocal buffer, start_para
        if not buffer:
            return
        section = " > ".join(heading_path) if heading_path else "body"
        chunks.extend(
            _chunk_text(
                "section",
                "\n".join(buffer),
                {"section": section, "paragraph": start_para},
                heading_path[:],
            )
        )
        buffer = []
        start_para = end_para + 1

    for index, paragraph in enumerate(doc.paragraphs, start=1):
        text = _normalize(paragraph.text)
        if not text:
            continue
        style_name = (paragraph.style.name if paragraph.style else "").lower()
        if style_name.startswith("heading"):
            flush(index - 1)
            level = _heading_level(style_name)
            heading_path[:] = heading_path[: level - 1]
            heading_path.append(text)
            sections += 1
            start_para = index + 1
            continue
        if not buffer:
            start_para = index
        buffer.append(text)
    flush(len(doc.paragraphs))
    if sections == 0 and chunks:
        sections = 1
    return _payload(os.path.basename(path), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", chunks, sections=sections)


def _heading_level(style_name: str) -> int:
    match = re.search(r"(\d+)", style_name)
    if not match:
        return 1
    return max(1, min(6, int(match.group(1))))


def _extract_pptx(path: str) -> dict:
    try:
        from pptx import Presentation
    except Exception:
        raise HTTPException(status_code=503, detail="pptx extractor unavailable")

    prs = Presentation(path)
    chunks: list[dict] = []
    slide_count = 0
    for index, slide in enumerate(prs.slides, start=1):
        slide_count += 1
        title = _slide_title(slide)
        lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = _normalize("".join(run.text for run in paragraph.runs))
                if text:
                    lines.append(text)
        section = title or f"slide {index}"
        heading_path = [title] if title else []
        chunks.extend(
            _chunk_text("slide", "\n".join(lines), {"slide": index, "section": section}, heading_path)
        )
    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    return _payload(os.path.basename(path), mime, chunks, sections=slide_count)


def _slide_title(slide: Any) -> str:
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        return _normalize(title_shape.text)
    return ""


def _extract_html(path: str) -> dict:
    result = _md.convert(path)
    chunks = _split_markdown_sections(result.text_content)
    return _payload(os.path.basename(path), "text/html", chunks, sections=len(chunks))


def _split_markdown_sections(markdown: str) -> list[dict]:
    chunks: list[dict] = []
    heading = ""
    buffer: list[str] = []

    def flush() -> None:
        nonlocal buffer
        if not buffer:
            return
        section = heading or "body"
        heading_path = [heading] if heading else []
        chunks.extend(_chunk_text("section", "\n".join(buffer), {"section": section}, heading_path))
        buffer = []

    for line in (markdown or "").splitlines():
        match = _atx_heading_re.match(line.strip())
        if match:
            flush()
            heading = _normalize(match.group(2))
            if heading:
                buffer.append(heading)
            continue
        buffer.append(line)
    flush()
    if not chunks:
        chunks = _chunk_text("section", markdown, {"section": "body"})
    return chunks


def _extract_csv(path: str) -> dict:
    import csv

    chunks: list[dict] = []
    lines: list[str] = []
    row_start = 0
    row_end = 0

    def flush() -> None:
        nonlocal lines, row_start, row_end
        if not lines:
            return
        chunks.extend(
            _chunk_text("rows", "\n".join(lines), {"row_start": row_start, "row_end": row_end})
        )
        lines = []
        row_start = 0
        row_end = 0

    with open(path, newline="", encoding="utf-8", errors="replace") as fh:
        for row_index, row in enumerate(csv.reader(fh), start=1):
            cells = [cell.strip() for cell in row if cell and cell.strip()]
            if not cells:
                continue
            if row_start == 0:
                row_start = row_index
            row_end = row_index
            lines.append(f"row {row_index}: " + " | ".join(cells))
            if len("\n".join(lines)) >= _max_chunk_chars or len(lines) >= 50:
                flush()
        flush()
    return _payload(os.path.basename(path), "text/csv", chunks)


def _extract_markdown(path: str) -> dict:
    result = _md.convert(path)
    chunks = _chunk_text("markdown", result.text_content, {})
    return _payload(os.path.basename(path), "", chunks)
