"""Generalist ETL: one schema for every file, dispatched by EXTENSION, never by content.

Three tables serve all formats, because every format decomposes the same way:

    source  -> one row per file
    unit    -> one row per sheet (.xlsx) / page (.pdf) / slide (.pptx) / file (.csv)
    datarow -> one row per spreadsheet row, or per unit of text

No column depends on what is inside a document, so a thousand unrelated spreadsheets
land without configuration. `unit.fingerprint` hashes the ordered header names + inferred
types, which is what later lets homogeneous sheets be clustered into real typed tables
(the D2 5107x case) while singletons stay generic.

DA RIFARE PRIMA DI INTEGRARLO — `resolve_header()` qui sotto reimplementa, senza test,
ciò che `internal/documents/filecard/` fa già e verifica:
`TestHeaderIsFoundBelowABannerRow`, `TestHeaderlessTableKeepsItsFirstRow`,
`TestCSVCardSniffsSemicolonsAndItalianDecimals`. Questo flow deve CHIAMARE filecard, non
duplicarlo. Vedi FINDINGS.md §5b. Il resto (dispatch per estensione, le tre tabelle,
il fingerprint di schema) resta valido.
"""

from __future__ import annotations

import csv
import dataclasses
import hashlib
import io
import json
import os
import pathlib
from collections.abc import AsyncIterator

import asyncpg
import cocoindex as coco
from cocoindex.connectors import localfs, postgres
from cocoindex.resources.file import PatternFilePathMatcher

DB_URL = os.environ["ETL_DATABASE_URL"]
PG_SCHEMA = os.environ.get("ETL_SCHEMA", "etl")
MAX_ROWS_PER_UNIT = int(os.environ.get("ETL_MAX_ROWS", "20000"))

PG_DB = coco.ContextKey[asyncpg.Pool]("pg_db")


@coco.lifespan
async def coco_lifespan(builder: coco.EnvironmentBuilder) -> AsyncIterator[None]:
    async with asyncpg.create_pool(DB_URL) as pool:
        async with pool.acquire() as con:
            await con.execute(f"CREATE SCHEMA IF NOT EXISTS {PG_SCHEMA}")
        builder.provide(PG_DB, pool)
        yield


@dataclasses.dataclass
class Source:
    id: str
    path: str
    ext: str
    size_bytes: int
    n_units: int


@dataclasses.dataclass
class Unit:
    id: str
    source_id: str
    ext: str
    unit_index: int
    unit_name: str
    n_rows: int
    n_cols: int
    header_json: str
    header_row: int      # -1 = no header found; rows are positional only
    header_score: float  # audit trail: how confident the heuristic was
    fingerprint: str


@dataclasses.dataclass
class DataRow:
    id: str
    unit_id: str
    row_num: int
    data_json: str   # {header_name: value} — empty {} when no header was found
    cells_json: str  # ["v0","v1",...] positional, ALWAYS written: a wrong header
                     # guess must never make a row unrecoverable


def _h(*parts: str) -> str:
    return hashlib.sha256("\x00".join(parts).encode("utf-8", "replace")).hexdigest()[:32]


def _kind(v) -> str:
    if v is None or v == "":
        return "null"
    if isinstance(v, bool):
        return "bool"
    if isinstance(v, int):
        return "int"
    if isinstance(v, float):
        return "float"
    s = str(v).strip()
    for cast, name in ((int, "int"), (float, "float")):
        try:
            cast(s.replace(",", "."))
            return name
        except ValueError:
            pass
    return "str"


def _fingerprint(header: list[str], sample: list[list]) -> str:
    """Ordered header names + the dominant type per column. Content-independent."""
    types = []
    for c in range(len(header)):
        kinds = [_kind(r[c]) for r in sample if c < len(r) and _kind(r[c]) != "null"]
        types.append(max(set(kinds), key=kinds.count) if kinds else "null")
    return _h("|".join(h.strip().lower() for h in header), "|".join(types))


def _profile(row: list) -> tuple[str, ...]:
    return tuple(_kind(v) for v in row)


def resolve_header(rows: list[list], scan: int = 12) -> tuple[int, float]:
    """Pick which row is the header, or return (-1, 0.0) for "there isn't one".

    Real spreadsheets put titles, blank spacers and merged banners above the header —
    measured on this corpus, 4 of 9 sheets did NOT have it on row 1. The score below is
    statistics about SHAPE only (fill ratio, all-text, uniqueness, and type divergence
    from the rows beneath). It never looks at what a value means, so it generalises to
    a corpus nobody has seen.
    """
    best, best_score = -1, 0.0
    width = max((len(r) for r in rows[:scan]), default=0)
    if not width:
        return -1, 0.0
    for i, row in enumerate(rows[:scan]):
        cells = [("" if v is None else str(v)).strip() for v in row]
        filled = sum(1 for c in cells if c)
        if filled < 2:
            continue
        fill = filled / width
        texty = sum(1 for c in cells if c and _kind(c) == "str") / max(filled, 1)
        uniq = len({c for c in cells if c}) / max(filled, 1)
        below = [r for r in rows[i + 1 : i + 1 + 20] if any(str(x or "").strip() for x in r)]
        if not below:
            continue
        # A header looks unlike its data: if the row below shares this row's type
        # profile, this is probably just another data row.
        diverge = sum(1 for r in below if _profile(r[:width]) != _profile(row[:width]))
        diverge /= len(below)
        score = 0.40 * fill + 0.25 * texty + 0.15 * uniq + 0.20 * diverge
        if score > best_score:
            best, best_score = i, score
    return (best, best_score) if best_score >= 0.55 else (-1, best_score)


# --- per-extension decomposition; each returns [(unit_name, rows, is_tabular)] ---
def _from_xlsx(content: bytes):
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    try:
        out = []
        for ws in wb.worksheets:
            rows = []
            for r in ws.iter_rows(values_only=True):
                rows.append(["" if c is None else str(c) for c in r])
                if len(rows) > MAX_ROWS_PER_UNIT:
                    break
            if rows:
                out.append((ws.title, rows, True))
        return out
    finally:
        wb.close()


def _from_csv(content: bytes, delim=","):
    text = content.decode("utf-8", "replace")
    rows = [r for r in csv.reader(io.StringIO(text), delimiter=delim)][:MAX_ROWS_PER_UNIT]
    return [("sheet1", rows, True)] if rows else []


def _from_pdf(content: bytes):
    import pypdfium2 as pdfium

    # Every pdfium handle must be closed explicitly and in creation order. Letting the
    # GC do it is what aborted the first run with `malloc_consolidate(): unaligned
    # fastbin chunk` at teardown — a native double-free, long after the data had
    # already been committed, so it looked like success until the exit code.
    doc = pdfium.PdfDocument(io.BytesIO(content))
    try:
        out = []
        for i in range(len(doc)):
            page = doc[i]
            tp = page.get_textpage()
            try:
                text = tp.get_text_range()
            finally:
                tp.close()
                page.close()
            out.append((f"page-{i + 1}", [[text]], False))
        return out
    finally:
        doc.close()


def _from_pptx(content: bytes):
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    out = []
    for i, slide in enumerate(prs.slides):
        chunks = [sh.text for sh in slide.shapes if getattr(sh, "has_text_frame", False)]
        out.append((f"slide-{i + 1}", [["\n".join(chunks)]], False))
    return out


def _from_text(content: bytes):
    return [("body", [[content.decode("utf-8", "replace")]], False)]


DISPATCH = {
    ".xlsx": _from_xlsx, ".xlsm": _from_xlsx,
    ".csv": _from_csv,
    ".tsv": lambda c: _from_csv(c, "\t"),
    ".pdf": _from_pdf,
    ".pptx": _from_pptx,
    ".txt": _from_text, ".md": _from_text, ".json": _from_text,
}


@coco.fn(memo=True)
async def process_file(
    file: localfs.File,
    source_t: postgres.TableTarget[Source],
    unit_t: postgres.TableTarget[Unit],
    row_t: postgres.TableTarget[DataRow],
) -> None:
    path = file.file_path.path.as_posix()
    ext = pathlib.Path(path).suffix.lower()
    handler = DISPATCH.get(ext)
    if handler is None:
        return  # unknown extension: skipped on purpose, never guessed at
    content = await file.read()
    try:
        units = handler(content)
    except Exception as e:  # a corrupt file must not stop the corpus
        print(f"  !! {path} ({ext}): {type(e).__name__}: {e}")
        return

    source_id = _h(path)
    source_t.declare_row(row=Source(
        id=source_id, path=path, ext=ext, size_bytes=len(content), n_units=len(units)))

    for idx, (unit_name, raw_rows, is_tabular) in enumerate(units):
        if is_tabular:
            hdr_i, score = resolve_header(raw_rows)
            header = ([str(h) for h in raw_rows[hdr_i]] if hdr_i >= 0 else [])
            data = raw_rows[hdr_i + 1 :] if hdr_i >= 0 else raw_rows
        else:
            hdr_i, score, header, data = 0, 1.0, ["text"], raw_rows

        unit_id = _h(source_id, unit_name, str(idx))
        width = max((len(r) for r in data[:50]), default=len(header))
        unit_t.declare_row(row=Unit(
            id=unit_id, source_id=source_id, ext=ext, unit_index=idx, unit_name=unit_name,
            n_rows=len(data), n_cols=max(width, len(header)),
            header_json=json.dumps(header, ensure_ascii=False),
            header_row=hdr_i, header_score=round(score, 3),
            fingerprint=_fingerprint(header or [f"c{i}" for i in range(width)], data[:25])))

        for n, r in enumerate(data):
            cells = ["" if v is None else str(v) for v in r]
            if not any(c.strip() for c in cells):
                continue  # blank spreadsheet rows are noise, not data
            named = {header[i]: c for i, c in enumerate(cells)
                     if i < len(header) and header[i].strip()} if header else {}
            row_t.declare_row(row=DataRow(
                id=_h(unit_id, str(n)), unit_id=unit_id, row_num=n,
                data_json=json.dumps(named, ensure_ascii=False),
                cells_json=json.dumps(cells, ensure_ascii=False)))


@coco.fn
async def app_main(sourcedir: pathlib.Path) -> None:
    source_t = await postgres.mount_table_target(
        PG_DB, table_name="source",
        table_schema=await postgres.TableSchema.from_class(Source, primary_key=["id"]),
        pg_schema_name=PG_SCHEMA)
    unit_t = await postgres.mount_table_target(
        PG_DB, table_name="unit",
        table_schema=await postgres.TableSchema.from_class(Unit, primary_key=["id"]),
        pg_schema_name=PG_SCHEMA)
    row_t = await postgres.mount_table_target(
        PG_DB, table_name="datarow",
        table_schema=await postgres.TableSchema.from_class(DataRow, primary_key=["id"]),
        pg_schema_name=PG_SCHEMA)

    files = localfs.walk_dir(
        sourcedir, recursive=True,
        path_matcher=PatternFilePathMatcher(
            included_patterns=[f"**/*{e}" for e in DISPATCH],
            excluded_patterns=["**/.*", "**/~$*"]))
    await coco.mount_each(process_file, files.items(), source_t, unit_t, row_t)


app = coco.App(
    coco.AppConfig(name="AuraEtlGeneric"),
    app_main,
    sourcedir=pathlib.Path(os.environ.get("CORPUS", "/corpus")),
)
